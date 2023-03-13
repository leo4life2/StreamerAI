from pptx.shapes.picture import Picture
from pptx import Presentation
import json

prs = Presentation('scripts/pres-old.pptx')

slide_objs = {}

# Shape ix to content:
# 3: Product selling points
# 4: Product name
# 5: Product size, retail price, sales price, shipping etc.
# 7: Brand description (same for all slides in pres-old.pptx, so emitted below)
# 8 or 9: Usage instructions, target audience, etc.

for i, slide in enumerate(prs.slides):
    slide_objs[i] = {}
    slide_objs[i]["name"] = slide.shapes[4].text
    slide_objs[i]["selling_points"] = slide.shapes[3].text
    slide_objs[i]["size_price"] = slide.shapes[5].text
    
    if isinstance(slide.shapes[8], Picture):
        slide_objs[i]["usage"] = slide.shapes[9].text
    else:
        slide_objs[i]["usage"] = slide.shapes[8].text
        
with open("scripts/pres-old-data.json", "w") as outfile:
    json.dump(slide_objs, outfile, ensure_ascii=False)